import os
import re
from pathlib import Path
from typing import Optional


KEYWORDS = {
    "University Docs": [
        "semester", "academic year", "credits", "course code",
        "roll number", "student name", "department", "university",
        "registrar", "faculty", "dean",
        "enrollment", "registration", "approved", "submitted",
        "internship approval", "student id", "application",
        "issued by", "official", "signature", "seal"
    ],
    "Technical Work": [
        "endpoint", "api", "pipeline", "build", "deployment",
        "server", "yaml", "json", "config", "docker", "kubernetes",
        "automate", "monitor", "logging", "debug", "ci/cd",
        "infrastructure", "version", "changelog",
        "developer", "engineer", "technical documentation",
        "readme", "implementation"
    ],
    "Capstone Work": [
        "abstract", "introduction", "methodology", "results",
        "evaluation", "discussion", "conclusion", "references",
        "proposal", "milestone", "phase", "final submission",
        "capstone project", "design", "implementation",
        "slides", "presentation", "data collection",
        "experiment", "analysis"
    ]
}


EXT_CONTENT_READERS = {}


def _text_normalize(s: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", s.lower())


def read_text_from_file(path: Path) -> str:
    """Return extracted text for supported extensions."""
    text = ""
    ext = path.suffix.lower()
    try:
        if ext == ".txt" or ext == ".md":
            text = path.read_text(encoding="utf-8", errors="ignore")
        elif ext == ".pdf":
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(str(path))
                parts = []
                for p in reader.pages:
                    try:
                        parts.append(p.extract_text() or "")
                    except Exception:
                        continue
                text = "\n".join(parts)
            except Exception:
                text = ""
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(str(path))
                text = "\n".join(p.text for p in doc.paragraphs)
            except Exception:
                text = ""
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                slides = []
                for slide in prs.slides:
                    for shp in slide.shapes:
                        if hasattr(shp, "text"):
                            slides.append(shp.text)
                text = "\n".join(slides)
            except Exception:
                text = ""
        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                parts = []
                for sheet in wb.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        for cell in row:
                            if cell is None:
                                continue
                            parts.append(str(cell))
                text = "\n".join(parts)
            except Exception:
                text = ""
        else:
            text = ""
    except Exception:
        text = ""
    return text


def classify_by_text(text: str) -> Optional[str]:
    s = _text_normalize(text)
    scores = {cat: 0 for cat in KEYWORDS}
    for cat, keywords in KEYWORDS.items():
        for kw in keywords:
            if kw in s:
                scores[cat] += 1
    best_cat = max(scores, key=lambda k: scores[k])
    if scores[best_cat] == 0:
        return None
    return best_cat


def classify_file(path: str) -> str:
    """Classify a single file path into one of the categories. Default fallback is 'Technical Work'"""
    p = Path(path)
    name_text = _text_normalize(p.stem + " " + p.name)
    # 1) filename
    for cat, keywords in KEYWORDS.items():
        for kw in keywords:
            if kw in name_text:
                return cat
    # 2) content
    content = read_text_from_file(p)
    if content:
        cat = classify_by_text(content)
        if cat:
            return cat

    # No file-extension based heuristics: rely only on filename and file content
    # for deterministic, explainable classification. Final fallback is Technical Work.
    return "Technical Work"


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python classifier.py <file_path>")
        sys.exit(1)
    print(classify_file(sys.argv[1]))
